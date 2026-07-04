"""PPTX Builder.

Builds .pptx files from structured slide data using python-pptx.
"""

import logging
import os
import tempfile
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# Color scheme
COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)  # Blue
COLOR_DARK = RGBColor(0x1F, 0x2D, 0x3D)  # Dark blue-gray
COLOR_LIGHT = RGBColor(0x6B, 0x72, 0x80)  # Gray
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BG = RGBColor(0xF3, 0xF4, 0xF6)  # Light gray bg


class PptxBuilder:
    """Build .pptx files from structured slide data."""

    def __init__(self, output_dir: str = ""):
        self._output_dir = output_dir or os.path.join(tempfile.gettempdir(), "training-ppts")
        os.makedirs(self._output_dir, exist_ok=True)

    def build(self, title: str, slides: list[dict]) -> str:
        """Build a .pptx file and return the local file path.

        Args:
            title: Presentation title (used for filename)
            slides: List of slide dicts with 'type' key

        Returns:
            Path to the generated .pptx file
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in slides:
            slide_type = slide_data.get("type", "content")
            if slide_type == "title":
                self._add_title_slide(prs, slide_data)
            elif slide_type == "content":
                self._add_content_slide(prs, slide_data)
            elif slide_type == "two_column":
                self._add_two_column_slide(prs, slide_data)
            elif slide_type == "summary":
                self._add_summary_slide(prs, slide_data)
            else:
                self._add_content_slide(prs, slide_data)

        safe_title = "".join(c for c in title if c.isalnum() or c in "_ -")[:50]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{safe_title}_{timestamp}.pptx"
        filepath = os.path.join(self._output_dir, filename)
        prs.save(filepath)

        logger.info(f"[PptxBuilder] Generated: {filepath}, slides={len(slides)}")
        return filepath

    def _add_title_slide(self, prs: Presentation, data: dict):
        """Title slide with centered title and subtitle."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_PRIMARY

        # Title
        title = data.get("title", "")
        subtitle = data.get("subtitle", "")

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.333), Inches(1))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = COLOR_WHITE
            p2.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs: Presentation, data: dict):
        """Content slide with title and bullet points."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title bar
        title_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = COLOR_PRIMARY
        title_shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.733), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        # Bullet points
        points = data.get("points", [])
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_DARK
            p.space_after = Pt(12)

    def _add_two_column_slide(self, prs: Presentation, data: dict):
        """Two-column layout slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title bar
        title_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = COLOR_PRIMARY
        title_shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.733), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        # Left column
        left_points = data.get("left", [])
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        for i, point in enumerate(left_points):
            if i == 0:
                p = tf_left.paragraphs[0]
            else:
                p = tf_left.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_DARK
            p.space_after = Pt(10)

        # Right column
        right_points = data.get("right", [])
        right_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.5), Inches(5))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        for i, point in enumerate(right_points):
            if i == 0:
                p = tf_right.paragraphs[0]
            else:
                p = tf_right.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_DARK
            p.space_after = Pt(10)

    def _add_summary_slide(self, prs: Presentation, data: dict):
        """Summary slide with a different visual style."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_PRIMARY

        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.333), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "总结")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

        # Summary points
        points = data.get("points", [])
        content_box = slide.shapes.add_textbox(Inches(2), Inches(2.2), Inches(9.333), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"✓ {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_WHITE
            p.space_after = Pt(14)
